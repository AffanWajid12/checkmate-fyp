from flask import Flask, request, jsonify
from flask_cors import CORS
import concurrent.futures
import hashlib
import pickle
import os
import requests
import re
from datetime import datetime
import io
import json
import math

# --- Dependency Imports ---
try:
    from sentence_transformers import SentenceTransformer
    import numpy as np
    import faiss
    import PyPDF2
    from pptx import Presentation
except ImportError:
    print("Dependencies not found. Please run: pip install -r requirements.txt")
    SentenceTransformer, np, faiss, PyPDF2, Presentation = None, None, None, None, None

# --- Flask App Initialization ---
app = Flask(__name__)
CORS(app)
CACHE_DIR = ".rag_cache"
RESOURCES_DIR = ".grading_resources"
for d in [CACHE_DIR, RESOURCES_DIR]:
    if not os.path.exists(d):
        os.makedirs(d)

class ResourceProcessor:
    @staticmethod
    def extract_text(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text
        elif ext == '.pptx':
            text = ""
            if not Presentation:
                raise ImportError("python-pptx not installed")
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        return None

class AIGrader:
    def __init__(self, assessment_id=None, strictness_level=None, additional_instructions=None):
        self.assessment_id = assessment_id
        self.strictness_level = strictness_level or ""
        self.additional_instructions = additional_instructions or ""
        self.llm_service_url = "http://localhost:5003/api/generate"
        self.faiss_index, self.text_chunks, self.rag_model = None, [], None
        
        # Load RAG model if available
        if faiss and SentenceTransformer:
            self.rag_model = SentenceTransformer('all-MiniLM-L6-v2')
            
            # Auto-load resources if assessment_id is provided
            if self.assessment_id:
                self._load_assessment_resources()
        else:
            print("--> AI Grader Initialized without RAG capabilities.")

    def _load_assessment_resources(self):
        """Loads and indexes resources from the .grading_resources/{assessment_id} directory."""
        res_dir = os.path.join(RESOURCES_DIR, str(self.assessment_id))
        if not os.path.exists(res_dir):
            return

        all_text = ""
        for filename in os.listdir(res_dir):
            file_path = os.path.join(res_dir, filename)
            try:
                extracted = ResourceProcessor.extract_text(file_path)
                if extracted:
                    all_text += f"\n\nSource: {filename}\n{extracted}"
            except Exception as e:
                print(f"Error extracting {filename}: {e}")

        if all_text.strip():
            print(f"--> Indexing resources for assessment {self.assessment_id}...")
            self._index_text(all_text)

    def _index_text(self, text_content):
        content_hash = hashlib.md5(text_content.encode()).hexdigest()
        cache_file = os.path.join(CACHE_DIR, f"{content_hash}.pkl")
        
        if os.path.exists(cache_file):
            with open(cache_file, "rb") as f:
                cached_data = pickle.load(f)
                self.text_chunks, self.faiss_index = cached_data['chunks'], faiss.deserialize_index(cached_data['index'])
            return
            
        # Basic chunking: split by paragraphs
        chunks = [p.strip() for p in text_content.split('\n\n') if p.strip()]
        if not chunks:
            return
            
        self.text_chunks = chunks
        embeddings = self.rag_model.encode(chunks, show_progress_bar=False)
        self.faiss_index = faiss.IndexFlatL2(embeddings.shape[1])
        self.faiss_index.add(embeddings)
        
        # Cache the index
        with open(cache_file, "wb") as f:
            pickle.dump({'chunks': self.text_chunks, 'index': faiss.serialize_index(self.faiss_index)}, f)

    def _retrieve_relevant_context(self, question, top_k=3):
        if not self.faiss_index or not self.rag_model:
            return ""
        question_embedding = self.rag_model.encode([question])
        _, indices = self.faiss_index.search(question_embedding, top_k)
        return "\n---\n".join([self.text_chunks[i] for i in indices[0]])

    def _grade_batch_with_ai(
        self,
        batch_items,
        additional_instructions=None,
        strictness_text=None
    ):
        """
        Grades a batch of questions in a single AI call.
        batch_items: list of (path_key, item_dict)
        """
        questions_block = ""
        for pk, item in batch_items:
            q_text = item.get('question', 'Missing question')
            ans_text = item.get('answer', '(No answer provided)')
            max_p = item.get('points', 0)
            rubric = item.get('rubric', 'No rubric provided')
            opts = item.get('options', [])
            
            # Robust extraction of option text if they are dicts
            clean_opts = []
            for o in opts:
                if isinstance(o, dict):
                    clean_opts.append(str(o.get('text', o)))
                else:
                    clean_opts.append(str(o))
            
            opts_str = f"Options: {', '.join(clean_opts)}" if clean_opts else ""
            
            # Re-integrate context per question in batch
            context = self._retrieve_relevant_context(q_text) if q_text else ""
            
            questions_block += f"""
### QUESTION ID: {pk}
Text: {q_text}
{opts_str}
Reference Context: {context}
Student Answer: {ans_text}
Max Marks: {max_p}
Rubric: {rubric}
---
"""

        teacher_block = f"Teacher's Instructions: {additional_instructions}\n" if additional_instructions else ""
        context_block = "Reference Material: (Use for factual verification if provided)\n"

        prompt = f"""
    You are an expert Computer Science teacher and strict grader.
    Evaluate the following student answers for the given questions.
    
    ### BATCH TO EVALUATE: (IMPORTANT: Grade each as a separate, isolated entity)
    {questions_block}

    ### GRADING CONFIG
    {teacher_block}
    Strictness Level: {strictness_text or "Balanced"}

    ---

    ### OUTPUT REQUIREMENTS
    - You must respond with a single, valid JSON object where keys are the QUESTION IDs.
    - Each value must be an object with score and categorized feedback.
    - DO NOT repeat or skip any IDs from the batch.
    - DO NOT mix feedback from one question into another.
    
    EXACT structure:
    {{
        "ID_HERE": {{
            "score": number (0 to Max Marks),
            "positive_points": ["point 1", "point 2"],
            "negative_points": ["point 1", "point 2"],
            "improvement_points": ["point 1", "point 2"],
            "summary": "concise summary"
        }}
    }}

    ### IMPORTANT:
    - Return ONLY the raw JSON object. 
    - DO NOT include markdown code blocks.
    - If the student's answer is blank or irrelevant, assign 0.
    """

        try:
            response_text = _call_ollama(prompt, timeout=180).strip()
            
            # Robust JSON cleaning
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}')
            if start_idx != -1 and end_idx != -1:
                json_str = response_text[start_idx:end_idx+1]
                batch_results = json.loads(json_str)
                
                # Post-process scores (strictness & bounds)
                final_batch = {}
                multiplier = self._strictness_multiplier(strictness_text or "")
                
                for pk, item in batch_items:
                    res = batch_results.get(pk, {})
                    raw_score = res.get('score', 0)
                    try: 
                        raw_score = float(raw_score)
                    except:
                        raw_score = 0
                    
                    max_p = item.get('points', 0)
                    adjusted = int(round(raw_score * multiplier))
                    final_batch[pk] = {
                        "score": max(0, min(adjusted, max_p)),
                        "positive_points": res.get('positive_points', []),
                        "negative_points": res.get('negative_points', []),
                        "improvement_points": res.get('improvement_points', []),
                        "summary": res.get('summary', response_text if not res else "")
                    }
                return final_batch
        except Exception as e:
            print(f"--> [BATCH-ERROR] {e}")
            
        # Fallback: empty results for this batch
        return {pk: {"score": 0, "positive_points": [], "negative_points": [], "improvement_points": [], "summary": "Batch grading failed."} for pk, _ in batch_items}


    def _grade_batch_with_llm(self, batch_items):
        """Standardizes interface for batch grading"""
        item_0 = batch_items[0][1]
        mode = "text"
        if "math" in item_0.get('answer_mode', ''): mode = "math"
        elif "coding" in item_0.get('answer_mode', ''): mode = "coding"
        
        mode_instructions = getattr(self, 'mode_instructions_map', {}).get(mode, "")
        
        print(f"--> [GRADING-BATCH] Starting batch of {len(batch_items)} items...")
        return self._grade_batch_with_ai(
            batch_items, 
            additional_instructions=mode_instructions,
            strictness_text=self.strictness_level
        )

    def grade_exam(self, student_exam):
        """
        Grade only leaf questions (questions without 'parts').
        Returns a mapping keyed by position-based IDs like "q-0-1-2" to avoid collisions
        and to ensure parent/container questions are NOT graded directly.
        """
        # Collect leaf tasks with their position path
        tasks = []  # list of (path_list, item)
        def walk(subparts, path_prefix, context_text=""):
            for idx, part in enumerate(subparts):
                cur_path = path_prefix + [idx]
                my_text = part.get('text', '') or part.get('question', '')
                
                # if there are sub-parts, recurse (container)
                if part.get('subparts'):
                    new_context = (context_text + "\n" + my_text).strip()
                    walk(part['subparts'], cur_path, new_context)
                else:
                    item = part.copy()
                    if context_text:
                        item['question'] = f"Context for this question:\n{context_text}\n\nQuestion:\n{my_text}".strip()
                    tasks.append((cur_path, item))
        walk(student_exam or [], [])

        results = {}
        # Batch processing
        BATCH_SIZE = 2
        for i in range(0, len(tasks), BATCH_SIZE):
            batch = tasks[i : i + BATCH_SIZE]
            # Convert tasks into batch_items expected by the batch grader
            batch_items = []
            for path, item in batch:
                path_key = "q-" + "-".join(map(str, path))
                batch_items.append((path_key, item))
            
            batch_results = self._grade_batch_with_llm(batch_items)
            
            for path_key, result in batch_results.items():
                # Fix variable shadowing in generator
                item = next((it for pk_internal, it in batch_items if pk_internal == path_key), None)
                if not item: continue
                
                results[path_key] = {
                    "score": result.get('score', 0),
                    "feedback": result,
                    "student_answer": item.get('answer', ''),
                    "max_score": int(item.get('points', 0)),
                    "question": item.get('question', ''),
                    "type": item.get('type', 'TextQuestion'),
                    "options": item.get('options', []),
                    "path": [int(x) for x in path_key.replace('q-', '').split('-')]
                }

        return results

    def _strictness_multiplier(self, strictness_text):
        """
        Map teacher's strictness instruction text to a deterministic multiplier.
        Multipliers >1 lenient (raise score), <1 strict (lower score).
        The values are conservative so we don't wildly change scores.
        """
        if not strictness_text:
            return 1.0
        t = strictness_text.lower()
        # priority map (try most specific phrases first)
        if "unforgiving" in t or "extremely strict" in t or "extremely unforgiving" in t:
            return 0.80
        if "extremely lenient" in t:
            return 1.15
        if "very strict" in t or "be very strict" in t:
            return 0.85
        if "very lenient" in t:
            return 1.10
        if "strict" in t and "slightly" in t:
            return 0.97
        if "slightly strict" in t:
            return 0.97
        if "slightly lenient" in t:
            return 1.03
        if "lenient" in t and "very" not in t and "extremely" not in t:
            return 1.05
        if "balanced" in t or "fair" in t or "be balanced" in t:
            return 1.0
        # fallback: look for "lenient" or "strict"
        if "lenient" in t:
            return 1.05
        if "strict" in t:
            return 0.95
        return 1.0


# ------------------ Helper to call Ollama ------------------
def _call_ollama(prompt, timeout=60, model="gemma", extra_kwargs=None):
    payload = {"prompt": prompt}
    if extra_kwargs:
        payload.update(extra_kwargs)
    resp = requests.post("http://localhost:5003/api/generate", json=payload, timeout=timeout)
    resp.raise_for_status()
    return resp.json().get("response", "")


def _generate_criteria_from_ai(question, num_columns):
    prompt = f"""
You are an expert teacher. Given the question below, produce a short list of rubric criterion NAMES.
Return ONLY a comma-separated list of criterion names.

Question: "{question}"
"""
    try:
        ai_text = _call_ollama(prompt, timeout=60)
        first_line = ai_text.splitlines()[0].strip()
        criteria = [c.strip() for c in first_line.split(",") if c.strip()]
        if not criteria:
            raise ValueError("AI returned no criteria.")
        return criteria
    except Exception as e:
        raise RuntimeError(f"Error generating criteria from AI: {e}")


def _calculate_mark_ranges(total_marks, num_columns):
    """
    Dynamically calculates mark ranges based on total marks and desired levels.
    Uses decimal precision for small total marks to avoid overlapping/zero ranges.
    """
    if num_columns <= 0:
        return []
    
    # Decide if we should use decimals (for low points) or integers
    use_decimals = total_marks <= 3
    precision = 1 if use_decimals else 0
    
    ranges = []
    names = {
        2: ["Excellent", "Needs Improvement"],
        3: ["Excellent", "Good", "Needs Improvement"],
        4: ["Excellent", "Good", "Satisfactory", "Needs Improvement"],
        5: ["Excellent", "Very Good", "Good", "Satisfactory", "Needs Improvement"],
    }.get(num_columns, [f"Level {i+1}" for i in range(num_columns)])
    
    step = total_marks / num_columns
    
    for i in range(num_columns):
        # Calculate bounds
        upper = total_marks - (i * step)
        lower = total_marks - ((i + 1) * step)
        
        # Ensure the last level always bottoms out at 0
        if i == num_columns - 1:
            lower = 0
            
        if use_decimals:
            u_str = f"{upper:.1f}"
            l_str = f"{lower:.1f}"
        else:
            u_str = f"{int(round(upper))}"
            l_str = f"{int(round(lower))}"
            
        # Avoid "1-1 Marks" if possible, but keep it if precision warrants
        if u_str == l_str:
            range_text = f"{u_str} Marks"
        else:
            range_text = f"{l_str}-{u_str} Marks"
            
        ranges.append(f"{range_text} ({names[i]})")
        
    return ranges

# -------------------------
# Rubric helper + bulk endpoints
# -------------------------

def _generate_rubric_matrix(question, total_marks, num_columns=3, provided_criteria=None, timeout_per_call=120):
    """
    Generate the entire rubric matrix in a single LLM call.
    Returns: [{"Criterion": "...", "Mark Range 1": "...", "Mark Range 2": "...", ...}, ...]
    """
    # 🧩 Normalize provided criteria into a list
    if isinstance(provided_criteria, str):
        provided_criteria = [c.strip() for c in provided_criteria.split(",") if c.strip()]

    criteria_instruction = ""
    if provided_criteria and isinstance(provided_criteria, list) and any(str(c).strip() for c in provided_criteria):
        criteria_instruction = f"Strictly use these specific criteria names: {', '.join(provided_criteria)}."
    else:
        criteria_instruction = "Automatically suggest 3-5 appropriate grading criteria (rows) based on the question."

    # Pre-calculate mark ranges to provide to the LLM as keys
    mark_ranges = _calculate_mark_ranges(total_marks, num_columns)
    mark_ranges_list_str = "\n".join([f"- \"{r}\"" for r in mark_ranges])

    # Construct a high-precision prompt for JSON output
    prompt = f"""
You are an expert teacher creating a comprehensive, cell-by-cell grading rubric for this question:
Question: "{question}"
Total Marks: {total_marks}

### Instructions:
1. Criteria (Rows): {criteria_instruction}
2. Performance Levels (Columns): You must generate descriptions for exactly these {num_columns} levels:
{mark_ranges_list_str}

3. Precision Rule: For low-mark questions (e.g., 1 or 2 points), ensure that the descriptions for high levels (like Excellent) specifically define what earns the full {total_marks} marks. DO NOT assign 0 marks for Excellent work.

4. Return the rubric ONLY as a valid JSON array of objects.
   - Each object must have a key "Criterion" for the name.
   - Each object must also have the following keys exactly as written for the level descriptions:
{mark_ranges_list_str}

### Example Format:
[
  {{
    "Criterion": "Content Accuracy",
    "{mark_ranges[0]}": "Details for top level...",
    "{mark_ranges[1]}": "Details for mid level...",
    ...
  }}
]

Return ONLY the raw JSON array. DO NOT include any markdown code blocks, preamble, or footer.
"""

    try:
        print(f"--> [RUBRIC-AI] Requesting single matrix for: '{question[:50]}...'")
        response_text = _call_ollama(prompt, timeout=timeout_per_call).strip()
        
        # Robust JSON cleaning: extract text within first '[' and last ']'
        start_idx = response_text.find('[')
        end_idx = response_text.rfind(']')
        if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
            json_str = response_text[start_idx:end_idx+1]
            rubric_matrix = json.loads(json_str)
            
            if isinstance(rubric_matrix, list):
                print(f"--> [RUBRIC-AI] Successfully generated matrix.")
                return rubric_matrix

        raise ValueError("Could not parse valid JSON array from LLM response.")

    except Exception as e:
        print(f"--> [RUBRIC-AI-ERROR] Failed: {e}")
        raise RuntimeError(f"Rubric generation failed: {e}")


def _generate_rubric_matrices_batch(questions_data, global_num_columns=3, timeout_per_call=180):
    """
    Optimized: Generate rubrics for multiple questions in a single LLM call.
    questions_data: list of {'id': str, 'question': str, 'points': int, 'columns': int, 'criteria': list}
    Returns: { question_id: rubric_matrix }
    """
    if not questions_data:
        return {}

    # Build a structured list for the prompt
    items_str = ""
    for q in questions_data:
        mark_ranges = _calculate_mark_ranges(q['points'], q['columns'] or global_num_columns)
        criteria_text = f"Criteria: {', '.join(q['criteria'])}" if q['criteria'] else "Suggest 3-4 criteria."
        items_str += f"""
ID: {q['id']}
Question: "{q['question']}"
Marks: {q['points']}
Columns: {', '.join(mark_ranges)}
{criteria_text}
---"""

    prompt = f"""
You are an expert teacher creating grading rubrics for multiple questions.
For each question below, generate a full rubric matrix.

### QUESTIONS:
{items_str}

### INSTRUCTIONS:
1. For each question, return a JSON object where the key is the ID and the value is a list of rubric rows.
2. Each rubric row must have a "Criterion" field and fields for each mark range level description provided above.
3. Return the response ONLY as a JSON object of the form: {{"ID_1": [{{...}}, {{...}}], "ID_2": [...]}}

Return ONLY the raw JSON object. No markdown, no preamble.
"""

    try:
        print(f"--> [BATCH-RUBRIC-AI] Requesting batch for {len(questions_data)} questions...")
        response_text = _call_ollama(prompt, timeout=timeout_per_call).strip()
        
        # Robust JSON cleaning
        start_idx = response_text.find('{')
        end_idx = response_text.rfind('}')
        if start_idx != -1 and end_idx != -1:
            json_str = response_text[start_idx:end_idx+1]
            batch_results = json.loads(json_str)
            return batch_results
            
        raise ValueError("No valid JSON object found in response.")
    except Exception as e:
        print(f"--> [BATCH-RUBRIC-AI-ERROR] Batch failed: {e}")
        return None  # Signal failure to top-level so it can retry individually if needed


def _generate_rubric_matrix_fallback(question, total_marks, num_columns, provided_criteria, timeout_per_call):
    # (Kept for compatibility, but the single prompt approach is preferred)
    pass



@app.route("/generate_question_rubric", methods=["POST"])
def generate_question_rubric():
    """Generate rubric for one question, respecting teacher criteria if given."""
    data = request.json or {}
    question = data.get("question")
    total_marks = data.get("points")
    num_columns = int(data.get("columns", 3))
    provided_criteria = data.get("criteria")
    path = data.get("path")

    if not question or total_marks is None:
        return jsonify({"error": "Missing question or marks"}), 400

    # 🧩 FIX: convert string to list if needed
    if isinstance(provided_criteria, str):
        provided_criteria = [c.strip() for c in provided_criteria.split(",") if c.strip()]

    try:
        rubric_matrix = _generate_rubric_matrix(
            question, total_marks, num_columns, provided_criteria
        )
    except Exception as e:
        return jsonify({"error": f"Failed to generate rubric: {e}"}), 500

    path_key = None
    if path is not None:
        path_key = "q-" + "-".join(map(str, path)) if isinstance(path, list) else str(path)

    payload = {"rubric": rubric_matrix}
    if path_key:
        payload["path_key"] = path_key
    return jsonify(payload), 200


@app.route("/generate_rubrics_bulk", methods=["POST"])
def generate_rubrics_bulk():
    """Generate rubrics for all leaf questions, using teacher criteria when provided.
    Supports per-question 'columns' field (q['columns']) as override of the global 'columns'.
    """
    data = request.json or {}
    student_exam = data.get("student_exam", [])
    global_num_columns = int(data.get("columns", 3))
    force = bool(data.get("force", False))
    timeout_per_call = int(data.get("timeout_per_call", 60))
    max_workers = int(data.get("max_workers", 6))

    # collect leaf questions with their position path
    leaves = []
    def walk(subparts, prefix):
        for idx, part in enumerate(subparts):
            cur_path = prefix + [idx]
            if part.get("subparts"):
                walk(part["subparts"], cur_path)
            else:
                leaves.append((cur_path, part))
    walk(student_exam or [], [])

    if not leaves:
        return jsonify({"generated": {}, "skipped": []}), 200

    generated = {}
    skipped = []

    # Batch questions together (e.g., 5 per call) to minimize API overhead
    BATCH_SIZE = 5
    to_generate = []
    for path, q in leaves:
        path_key = "q-" + "-".join(map(str, path))
        if q.get("rubric") and not force:
            skipped.append(path_key)
            continue
        
        # Support both 'points' (legacy/single) and 'total_marks' (standard blueprint)
        raw_points = q.get("total_marks") if q.get("total_marks") is not None else q.get("points")
        try:
            leaf_points = float(raw_points) if raw_points is not None else 0.0
        except (ValueError, TypeError):
            leaf_points = 0.0

        leaf_columns = global_num_columns
        try:
            if q.get("rubric_columns") is not None:
                leaf_columns = int(q.get("rubric_columns"))
            elif q.get("columns") is not None:
                leaf_columns = int(q.get("columns"))
        except:
            pass

        provided_criteria = q.get("criteria") or q.get("suggested_criteria")
        if isinstance(provided_criteria, str):
            provided_criteria = [c.strip() for c in provided_criteria.split(",") if c.strip()]

        to_generate.append({
            'id': path_key,
            'question': q.get("question") or q.get("text", ""),
            'points': leaf_points,
            'columns': int(leaf_columns),
            'criteria': provided_criteria
        })

    if not to_generate:
        return jsonify({"generated": {}, "skipped": skipped}), 200

    # Process in batches
    for i in range(0, len(to_generate), BATCH_SIZE):
        batch = to_generate[i:i + BATCH_SIZE]
        results = _generate_rubric_matrices_batch(batch, global_num_columns, timeout_per_call * 2)
        
        if results and isinstance(results, dict):
            for pk, rubric in results.items():
                generated[pk] = rubric
        else:
            # If batch fails, try each question in this batch individually
            print(f"--> [BATCH-FAILED] Retrying batch {i//BATCH_SIZE + 1} individually...")
            for item in batch:
                try:
                    rubric = _generate_rubric_matrix(
                        item['question'], item['points'], item['columns'], 
                        provided_criteria=item['criteria'], 
                        timeout_per_call=timeout_per_call
                    )
                    generated[item['id']] = rubric
                except Exception as e:
                    generated[item['id']] = {"error": str(e)}

    return jsonify({"generated": generated, "skipped": skipped}), 200



@app.route('/delete_rubric', methods=['POST'])
def delete_rubric():
    """
    Delete a rubric for a specific question. Accepts:
      - { "path_key": "q-0-1" }
      - or { "path": [0,1] }

    Returns { "deleted": "<path_key>" } or 404 if not found.
    Note: this endpoint does NOT mutate your original exam storage — it returns the path_key so frontend
    can update its `exam_questions` state. If you want server-side persistence, adapt to your DB.
    """
    data = request.json or {}
    path_key = data.get('path_key')
    path = data.get('path')

    if not path_key and path is None:
        return jsonify({"error": "Provide 'path_key' or 'path'."}), 400

    if path_key is None and isinstance(path, list):
        path_key = "q-" + "-".join(map(str, path))

    # In this simple implementation we just return the path_key so the frontend can clear it locally.
    # If you have server-side exam storage, delete the rubric there and return a status.
    return jsonify({"deleted": path_key}), 200



@app.route('/upload_grading_resource', methods=['POST'])
def upload_grading_resource():
    """Endpoint for uploading TXT, PDF, or PPTX resources for RAG."""
    assessment_id = request.form.get('assessment_id')
    if not assessment_id:
        return jsonify({"error": "Missing assessment_id"}), 400

    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Empty filename"}), 400

    # Create assessment dir
    res_dir = os.path.join(RESOURCES_DIR, assessment_id)
    if not os.path.exists(res_dir):
        os.makedirs(res_dir)

    file_path = os.path.join(res_dir, file.filename)
    file.save(file_path)

    # Try extracting to verify support
    try:
        text = ResourceProcessor.extract_text(file_path)
        if not text:
            return jsonify({"error": "Could not extract text from file"}), 400
        
        return jsonify({
            "message": f"Resource {file.filename} uploaded and processed",
            "filename": file.filename,
            "status": "success"
        }), 200
    except Exception as e:
        if os.path.exists(file_path):
            os.remove(file_path)
        return jsonify({"error": f"Failed to process resource: {str(e)}"}), 500

@app.route('/clear_grading_resources', methods=['POST'])
def clear_grading_resources():
    assessment_id = request.json.get('assessment_id')
    if not assessment_id:
        return jsonify({"error": "Missing assessment_id"}), 400

    res_dir = os.path.join(RESOURCES_DIR, assessment_id)
    if os.path.exists(res_dir):
        import shutil
        shutil.rmtree(res_dir)
        return jsonify({"message": f"Resources for assessment {assessment_id} cleared"}), 200
    
    return jsonify({"message": "No resources found to clear"}), 200

@app.route('/get_grading_resources', methods=['GET'])
def get_grading_resources():
    assessment_id = request.args.get('assessment_id')
    print(f"--> [GET_RESOURCES] Request for Assessment ID: {assessment_id}")
    if not assessment_id:
        return jsonify([])

    res_dir = os.path.join(RESOURCES_DIR, assessment_id)
    if not os.path.exists(res_dir):
        print(f"--> [GET_RESOURCES] Directory not found: {res_dir}")
        return jsonify([])

    try:
        files = [f for f in os.listdir(res_dir) if os.path.isfile(os.path.join(res_dir, f))]
        print(f"--> [GET_RESOURCES] Found files: {files}")
        return jsonify(files)
    except Exception as e:
        print(f"--> [GET_RESOURCES] Error listing directory: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/grade', methods=['POST'])
def grade():
    print(f"\n{'=' * 50}\n--- Received GRADING request ---\n{'=' * 50}")
    data = request.json
    try:
        grader = AIGrader(
            assessment_id=data.get('assessment_id'),
            strictness_level=data.get('strictness_level'),
            additional_instructions=data.get('additional_instructions', '')
        )
        # Store the mode-specific instructions map
        grader.mode_instructions_map = {
            "text": data.get("text_instructions", ""),
            "math": data.get("math_instructions", ""),
            "coding": data.get("coding_instructions", "")
        }
        results = grader.grade_exam(data['student_exam'])
        return jsonify(results)
    except Exception as e:
        print(f"--> [FATAL ERROR] {e}")
        return jsonify({"error": str(e)}), 500



if __name__ == '__main__':
    print("Starting Flask AI Grading Server on Port 5004 with Teacher-Guided Rubrics...")
    app.run(host='0.0.0.0', port=5004)
